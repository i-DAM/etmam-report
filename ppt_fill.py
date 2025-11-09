# -*- coding: utf-8 -*-
from __future__ import annotations

from io import BytesIO
import re

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor


def _norm_ar(s: str) -> str:
    if not isinstance(s, str):
        return ""
    s = s.strip()
    s = (
        s.replace("أ", "ا")
         .replace("إ", "ا")
         .replace("آ", "ا")
         .replace("ى", "ي")
         .replace("ة", "ه")
    )
    s = re.sub(r"\s+", " ", s)
    s = s.replace(" -", "-").replace("- ", "-")
    return s


def _keywords(s: str) -> set[str]:
    s = _norm_ar(s)
    stop = {
        "الاداره", "اداره", "الادارة",
        "العامه", "عامه", "العامة",
        "بلديه", "بلدية",
        "امانه", "امانة",
        "مدينه", "مدينة", "المدينه",
        "منطقه", "منطقة", "المنطقه",
    }
    toks: list[str] = []
    for t in s.split():
        if not t:
            continue
        if t.startswith("لل") and len(t) > 2:
            t = t[2:]
        elif t.startswith("ال") and len(t) > 2:
            t = t[2:]
        elif t.startswith("ل") and len(t) > 2:
            t = t[1:]
        if t and t not in stop:
            toks.append(t)
    return set(toks)


def _find_main_table(slide):
    tables = [sh.table for sh in slide.shapes if sh.has_table]
    if not tables:
        raise RuntimeError("لم يتم العثور على جدول في الشريحة.")
    main = max(tables, key=lambda t: len(t.columns))
    return main


def _detect_main_columns(table):
    header = table.rows[0].cells
    m = dict(admin=None, open=None, reopen=None, near=None, late=None, other=None)

    for i, c in enumerate(header):
        txt = c.text
        tnorm = _norm_ar(txt)

        if any(k in tnorm for k in ["الادارات", "الاداره", "الادارة"]):
            m["admin"] = i
        if "البلاغات المفتوحه" in tnorm or "البلاغات المفتوحة" in txt:
            m["open"] = i
        if "اعادة فتح" in tnorm or "المعاد فتحها" in tnorm:
            m["reopen"] = i
        if "قارب" in tnorm or "قارب على تجاوز" in tnorm:
            m["near"] = i
        if "المتأخره" in tnorm or "المتأخرة" in txt or "تجاوز sla" in tnorm:
            m["late"] = i
        if "مصادر اخرى" in tnorm or "مصادر أخرى" in txt:
            m["other"] = i

    if m["admin"] is None:
        m["admin"] = 0

    return m


def _make_shape_text_white(shape) -> None:
    if not shape.has_text_frame:
        return
    for p in shape.text_frame.paragraphs:
        for run in p.runs:
            run.font.color.rgb = RGBColor(255, 255, 255)


def _make_cell_text_white(cell) -> None:
    if not cell.text_frame:
        return
    for p in cell.text_frame.paragraphs:
        for run in p.runs:
            run.font.color.rgb = RGBColor(255, 255, 255)


def _fill_date_placeholder(prs: Presentation, closed_dt: pd.Timestamp) -> None:
    weekday_map = {
        0: "الاثنين",
        1: "الثلاثاء",
        2: "الأربعاء",
        3: "الخميس",
        4: "الجمعة",
        5: "السبت",
        6: "الأحد",
    }
    weekday_ar = weekday_map[closed_dt.weekday()]
    date_str = closed_dt.strftime("%Y-%m-%d")
    full_text = f"التحديث اليومي {weekday_ar} {date_str}"

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            txt = shape.text
            if "{{DATE}}" not in txt and "{DATE}" not in txt:
                continue
            txt = txt.replace("{{DATE}}", full_text).replace("{DATE}", full_text)
            shape.text = txt
            _make_shape_text_white(shape)


def _fill_total_placeholders(
    prs: Presentation,
    open_total_all: int,
    near_total_all: int,
    late_total_all: int,
    other_total_all: int,
) -> None:
    repl = {
        "{OPEN_TOTAL}": str(open_total_all),
        "{{OPEN_TOTAL}}": str(open_total_all),
        "{NEAR_SLA_TOTAL}": str(near_total_all),
        "{{NEAR_SLA_TOTAL}}": str(near_total_all),
        "{LATE_TOTAL}": str(late_total_all),
        "{{LATE_TOTAL}}": str(late_total_all),
        "{OTHER_TOTAL}": str(other_total_all),
        "{{OTHER_TOTAL}}": str(other_total_all),
    }

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text
            changed = False
            for token, val in repl.items():
                if token in text:
                    text = text.replace(token, val)
                    changed = True
            if changed:
                shape.text = text


def _find_reopen_column(df: pd.DataFrame) -> str | None:
    for col in df.columns:
        s = _norm_ar(str(col))
        if ("معلق" in s or "مغلق" in s) and ("فتح" in s):
            return col
    return None


def _find_sla_columns(df: pd.DataFrame):
    near_col = None
    late_col = None
    for col in df.columns:
        s = _norm_ar(str(col))
        if "قارب" in s:
            near_col = col
        if "تجاوز" in s:
            late_col = col
    return near_col, late_col


def _match_admin_indices(admin_in_slide: str, all_admins: list[str]) -> list[str]:
    s = _norm_ar(admin_in_slide)
    kw_s = _keywords(s)

    matches: list[str] = []

    is_muni = ("بلديه" in s or "بلدية" in s)
    is_south = "جنوب" in s
    is_north = "شمال" in s
    is_center = "وسط" in s

    is_main_clean = (
        ("النظافه" in s or "النظافة" in s)
        and ("اداره" in s or "الاداره" in s or "الادارة" in s)
        and ("عامه" in s or "العامه" in s or "العامة" in s)
    )

    is_main_projects = (
        ("مشاريع" in s)
        and ("اداره" in s or "الاداره" in s or "الادارة" in s)
        and ("عامه" in s or "العامه" in s or "العامة" in s)
    )

    for name in all_admins:
        n = _norm_ar(name)
        kw_n = _keywords(n)

        if is_muni and (is_south or is_north or is_center):
            if is_south and "جنوب" in n and (
                "تعمير" in n or ("رقابه" in n and "الخدمات" in n)
            ):
                matches.append(name)
                continue

            if is_north and "شمال" in n and (
                "تعمير" in n or ("رقابه" in n and "الخدمات" in n)
            ):
                matches.append(name)
                continue

            if is_center and "وسط" in n and (
                "تعمير" in n or ("رقابه" in n and "الخدمات" in n)
            ):
                matches.append(name)
                continue

            continue

        if is_main_clean:
            if "النظافه" in n or "النظافة" in n:
                matches.append(name)
                continue

        if is_main_projects:
            if "مشاريع" in n:
                matches.append(name)
                continue

        if ("تشغيل" in s and "صيان" in s):
            if ("تشغيل" in n and "صيان" in n) or ("اناره" in n) or ("انارة" in n):
                matches.append(name)
                continue

        if ("الاصحاح" in s and "البيئي" in s):
            if "الاصحاح" in n and "البيئي" in n:
                matches.append(name)
                continue

        if n == s:
            matches.append(name)
            continue

        if kw_n == kw_s or kw_n.issubset(kw_s) or kw_s.issubset(kw_n):
            matches.append(name)
            continue

    seen = set()
    ordered = []
    for name in all_admins:
        if name in matches and name not in seen:
            ordered.append(name)
            seen.add(name)
    return ordered


def _aggregate_for_admin(admin_in_slide: str, df: pd.DataFrame) -> pd.Series:
    if df is None or df.empty:
        return pd.Series(0, index=df.columns if hasattr(df, "columns") else [])

    df = df.copy()
    df.index = df.index.astype(str)
    if "الإجمالي الكلي" in df.index:
        df = df.drop(index="الإجمالي الكلي")

    all_admins = list(df.index)
    idxs = _match_admin_indices(admin_in_slide, all_admins)
    if not idxs:
        return pd.Series(0, index=df.columns)

    sub = df.loc[idxs]
    if isinstance(sub, pd.Series):
        return sub
    return sub.sum(axis=0)


def _fill_main_table(
    table,
    p_open: pd.DataFrame,
    p_sla: pd.DataFrame,
    p_other: pd.DataFrame | None,
):
    colmap = _detect_main_columns(table)

    reopen_col = _find_reopen_column(p_open)
    near_col, late_col = _find_sla_columns(p_sla)

    total_open = total_reopen = total_near = total_late = total_other = 0
    total_row_idx = None

    n_rows = len(table.rows)

    for r in range(1, n_rows):
        cells = table.rows[r].cells
        admin_txt = cells[colmap["admin"]].text.strip()
        norm_admin = _norm_ar(admin_txt)

        if not admin_txt:
            continue

        if "الاجمالي" in norm_admin:
            total_row_idx = r
            continue

        vals_open = _aggregate_for_admin(admin_txt, p_open)
        if reopen_col is not None:
            reopen_val = int(vals_open.get(reopen_col, 0))
        else:
            reopen_val = 0
        open_val = int(vals_open.sum() - reopen_val)

        vals_sla = _aggregate_for_admin(admin_txt, p_sla)
        near_val = int(vals_sla.get(near_col, 0)) if near_col else 0
        late_val = int(vals_sla.get(late_col, 0)) if late_col else 0

        if p_other is not None and not p_other.empty:
            vals_other = _aggregate_for_admin(admin_txt, p_other)
            other_val = int(vals_other.sum())
        else:
            other_val = 0

        if colmap["open"] is not None:
            cells[colmap["open"]].text = str(open_val)
        if colmap["reopen"] is not None:
            cells[colmap["reopen"]].text = str(reopen_val)
        if colmap["near"] is not None:
            cells[colmap["near"]].text = str(near_val)
        if colmap["late"] is not None:
            cells[colmap["late"]].text = str(late_val)
        if colmap["other"] is not None:
            cells[colmap["other"]].text = str(other_val)

        total_open += open_val
        total_reopen += reopen_val
        total_near += near_val
        total_late += late_val
        total_other += other_val

    if total_row_idx is not None:
        cells = table.rows[total_row_idx].cells
        if colmap["open"] is not None:
            cells[colmap["open"]].text = str(total_open)
            _make_cell_text_white(cells[colmap["open"]])
        if colmap["reopen"] is not None:
            cells[colmap["reopen"]].text = str(total_reopen)
            _make_cell_text_white(cells[colmap["reopen"]])
        if colmap["near"] is not None:
            cells[colmap["near"]].text = str(total_near)
            _make_cell_text_white(cells[colmap["near"]])
        if colmap["late"] is not None:
            cells[colmap["late"]].text = str(total_late)
            _make_cell_text_white(cells[colmap["late"]])
        if colmap["other"] is not None:
            cells[colmap["other"]].text = str(total_other)
            _make_cell_text_white(cells[colmap["other"]])

    return total_open, total_near, total_late, total_other


def _fill_side_tables(
    slide,
    main_table,
    p_open: pd.DataFrame,
    p_sla: pd.DataFrame,
    p_other: pd.DataFrame | None,
):
    reopen_col = _find_reopen_column(p_open)
    near_col, late_col = _find_sla_columns(p_sla)

    for shape in slide.shapes:
        if not shape.has_table:
            continue
        table = shape.table
        if table is main_table:
            continue

        colmap = _detect_main_columns(table)
        admin_idx = colmap.get("admin")
        if admin_idx is None:
            continue

        n_rows = len(table.rows)
        for r in range(1, n_rows):
            cells = table.rows[r].cells
            admin_txt = cells[admin_idx].text.strip()
            if not admin_txt:
                continue
            if "الاجمالي" in _norm_ar(admin_txt):
                continue

            vals_open = _aggregate_for_admin(admin_txt, p_open)
            if reopen_col is not None:
                reopen_val = int(vals_open.get(reopen_col, 0))
            else:
                reopen_val = 0
            open_val = int(vals_open.sum() - reopen_val)

            vals_sla = _aggregate_for_admin(admin_txt, p_sla)
            near_val = int(vals_sla.get(near_col, 0)) if near_col else 0
            late_val = int(vals_sla.get(late_col, 0)) if late_col else 0

            if p_other is not None and not p_other.empty:
                vals_other = _aggregate_for_admin(admin_txt, p_other)
                other_val = int(vals_other.sum())
            else:
                other_val = 0

            if colmap["open"] is not None:
                cells[colmap["open"]].text = str(open_val)
            if colmap["reopen"] is not None:
                cells[colmap["reopen"]].text = str(reopen_val)
            if colmap["near"] is not None:
                cells[colmap["near"]].text = str(near_val)
            if colmap["late"] is not None:
                cells[colmap["late"]].text = str(late_val)
            if colmap["other"] is not None:
                cells[colmap["other"]].text = str(other_val)


_CARD_PATTERN = re.compile(r"\{CARD_(OPEN|NEAR|LATE|OTHER)([^}]*)\}", re.DOTALL)


def _fill_left_cards(
    prs: Presentation,
    p_open: pd.DataFrame,
    p_sla: pd.DataFrame,
    p_other: pd.DataFrame | None,
):
    reopen_col = _find_reopen_column(p_open)
    near_col, late_col = _find_sla_columns(p_sla)

    cache: dict[str, dict[str, int]] = {}

    def get_metrics(admin_name: str) -> dict[str, int]:
        admin_name = admin_name.strip()
        if not admin_name:
            return {"open": 0, "near": 0, "late": 0, "other": 0}

        if admin_name in cache:
            return cache[admin_name]

        vals_open = _aggregate_for_admin(admin_name, p_open)
        if reopen_col is not None:
            reopen_val = int(vals_open.get(reopen_col, 0))
        else:
            reopen_val = 0
        open_val = int(vals_open.sum() - reopen_val)

        vals_sla = _aggregate_for_admin(admin_name, p_sla)
        near_val = int(vals_sla.get(near_col, 0)) if near_col else 0
        late_val = int(vals_sla.get(late_col, 0)) if late_col else 0

        if p_other is not None and not p_other.empty:
            vals_other = _aggregate_for_admin(admin_name, p_other)
            other_val = int(vals_other.sum())
        else:
            other_val = 0

        metrics = {
            "open": open_val,
            "near": near_val,
            "late": late_val,
            "other": other_val,
        }
        cache[admin_name] = metrics
        return metrics

    from pptx.enum.text import PP_ALIGN

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            txt = shape.text
            if "CARD_" not in txt:
                continue

            def repl(match: re.Match) -> str:
                kind = match.group(1).upper()
                rest = match.group(2) or ""

                admin = rest.replace("\n", " ")
                admin = admin.replace(":", " ").replace("ـ", " ").replace("-", " ")
                admin = admin.strip()

                metrics = get_metrics(admin)
                key_map = {
                    "OPEN": "open",
                    "NEAR": "near",
                    "LATE": "late",
                    "OTHER": "other",
                }
                return str(metrics.get(key_map[kind], 0))

            new_txt = _CARD_PATTERN.sub(repl, txt)
            if new_txt != txt:
                shape.text = new_txt
                tf = shape.text_frame
                for p in tf.paragraphs:
                    p.alignment = PP_ALIGN.CENTER
                    for run in p.runs:
                        run.font.color.rgb = RGBColor(255, 255, 255)


def fill_ppt(
    template_path: str,
    p_open: pd.DataFrame,
    p_sla: pd.DataFrame,
    p_other: pd.DataFrame | None,
    closed_dt: pd.Timestamp,
) -> BytesIO:
    prs = Presentation(template_path)
    slide = prs.slides[0]

    main_table = _find_main_table(slide)

    open_total_all, near_total_all, late_total_all, other_total_all = _fill_main_table(
        main_table, p_open, p_sla, p_other
    )

    _fill_side_tables(slide, main_table, p_open, p_sla, p_other)

    _fill_left_cards(prs, p_open, p_sla, p_other)

    _fill_date_placeholder(prs, closed_dt)

    _fill_total_placeholders(
        prs,
        open_total_all=open_total_all,
        near_total_all=near_total_all,
        late_total_all=late_total_all,
        other_total_all=other_total_all,
    )

    out = BytesIO()
    prs.save(out)
    out.seek(0)
    return out
